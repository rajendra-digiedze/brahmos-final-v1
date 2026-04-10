from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

def add_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        p = content.add_paragraph()
        if isinstance(point, tuple):
            text, level = point
            p.text = text
            p.level = level
        else:
            p.text = point
            if i == 0:
                content.text = point # First item replaces the default text
                continue

# Slide 1: Title Slide (1)
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "dZshield Enterprise SOC"
subtitle.text = "Automated Log Analysis & Threat Mitigation\nEnd-to-End Integration Flow\n\n10-Page Architecture Breakdown"

# Slide 2: Introduction (2)
add_slide(prs, "1. Introduction to dZshield SOC", [
    "dZshield is a next-generation Security Operations Center (SOC) dashboard.",
    ("It focuses on Zero-Trust Network Topology and automated threat hunting.", 1),
    ("Monitors network traffic, endpoint data, and system logs in real-time.", 1),
    "The primary goal is minimizing manual analyst intervention through AI.",
    ("Automates triage, context retrieval, and incident response operations.", 1)
])

# Slide 3: Core Technology Stack (3)
add_slide(prs, "2. Core Technology Stack", [
    "The system relies on four primary pillars of technology:",
    ("1. Nify: High-speed log ingestion and telemetry routing.", 1),
    ("2. n8n: Complex workflow orchestration and event triggers.", 1),
    ("3. Qdrant: High-performance Vector Database for historical context.", 1),
    ("4. Llama 3 (AI): Large Language Model for intelligent investigation.", 1)
])

# Slide 4: Component 1 - Nify (4)
add_slide(prs, "3. Log Ingestion (Nify)", [
    "Nify serves as the central nervous system for data collection.",
    ("Ingests JSON-formatted logs from web servers, firewalls, and endpoints.", 1),
    ("Normalizes data into a standard schema for downstream processing.", 1),
    ("Ensures zero packet loss during high-traffic cyber attacks.", 1),
    "Acts as the first step in the dZshield pipeline."
])

# Slide 5: Component 2 - n8n Orchestrator (5)
add_slide(prs, "4. Workflow Automation (n8n)", [
    "n8n replaces traditional hard-coded SOC playbooks.",
    ("Uses a visual, node-based routing system to connect APIs.", 1),
    ("Listens for incoming alerts from Nify automatically.", 1),
    ("Triggers the AI agent to start processing when critical logs appear.", 1),
    ("Responsible for sending outbound notifications (Email, Webhooks).", 1)
])

# Slide 6: Component 3 - Qdrant Vector DB (6)
add_slide(prs, "5. Threat Memory Context (Qdrant)", [
    "Qdrant powers the system's Long-Term Context and Memory.",
    ("Stores logs as high-dimensional vector embeddings.", 1),
    ("Allows the system to perform Semantic Searches instead of keyword searches.", 1),
    ("Helps identify similar past attacks to determine current threat severity.", 1),
    "Retrieval-Augmented Generation (RAG) significantly reduces AI hallucinations."
])

# Slide 7: Component 4 - AI Analysis (7)
add_slide(prs, "6. Intelligent Triage (Llama 3 AI)", [
    "Llama 3 replaces the human Level-1 SOC Analyst.",
    ("Powered by LangGraph to process multi-step reasoning workflows.", 1),
    ("Reads the incoming log and the historical context provided by Qdrant.", 1),
    ("Generates a plain-English explanation of the attack vector.", 1),
    ("Recommends actionable mitigation steps immediately.", 1)
])

# Slide 8: The End-to-End Workflow (8)
add_slide(prs, "7. The Automated Incident Workflow", [
    "How an attack is handled from start to finish:",
    ("Detection: Nify detects an anomalous login spike.", 1),
    ("Forwarding: Log is pushed to n8n Webhook.", 1),
    ("Contextualization: n8n queries Qdrant for similar historical logins.", 1),
    ("Analysis: Llama 3 reviews the data and flags it as a 'Brute Force Attack'.", 1),
    ("Response: n8n sends a critical alert to the admin and updates the UI.", 1)
])

# Slide 9: User Interface & Dashboard (9)
add_slide(prs, "8. The SOC Dashboard Interface", [
    "The React-based Frontend provides total visibility into the automated operations.",
    ("Interactive Chat4ED Agent: Query logs naturally (e.g. 'Show me errors from 5 mins ago').", 1),
    ("Dynamic Visualizations: Auto-syncing charts reflecting data across custom timeframes.", 1),
    ("Integrated UI Tabs: Full embedded access to n8n and Qdrant interfaces.", 1),
    ("Report Generation: Export massive datasets and graphs directly to Excel/PDF.", 1)
])

# Slide 10: Conclusion (10)
add_slide(prs, "9. Conclusion & Capabilities", [
    "dZshield fundamentally accelerates incident response times.",
    ("Combines automation (n8n), context (Qdrant), and intelligence (Llama).", 1),
    ("Scalable architecture that bypasses standard container limits when deployed natively.", 1),
    ("Fully customizable flow rules allow the system to adapt to new zero-day threats.", 1),
    "End Result: A robust, enterprise-grade, Zero-Trust environment."
])

prs.save('dZshield_SOC_Detailed_Flow.pptx')
print("10-PAGE PPT GENERATED SUCCESSFULLY.")

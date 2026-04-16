import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_brahmos_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Define slide layouts (0 = Title, 1 = Title & Content)
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    
    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "BrahMos Enterprise SOC Dashboard"
    subtitle.text = "7-Layer Artificial Intelligence Architecture\nAutomated Cyber-Defense Workflow"
    
    # Stylize Title
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # ----------------------------------------------------
    # Slide 2: High-Level Data Pipeline
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    body = shapes.placeholders[1]
    
    title.text = "End-to-End Packet Pipeline Overview"
    tf = body.text_frame
    tf.text = "How a packet is intercepted and analyzed natively:"
    
    p = tf.add_paragraph()
    p.text = "1. Interception: Edge firewalls/Windows local logs emit raw syslog data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Routing: Data enters the FastAPI backend and is stored into ACID-compliant MySQL."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Embeddings: High severity alerts are embedded mathematically and pushed to Qdrant Vector Cloud."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Analytics: Frontend React application polls the MySQL layers dynamically to draw Area and Pie charts."
    p.level = 1

    # ----------------------------------------------------
    # Slide 3: L1 & L2 Details
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    body = shapes.placeholders[1]
    
    title.text = "Layer 1 & Layer 2: Foundations"
    tf = body.text_frame
    
    tf.text = "Layer 1: Infrastructure (Network Hardware/Virtual)"
    p = tf.add_paragraph()
    p.text = "Tools: Kubernetes, dZ GaaS, Fortigate Firewalls"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: Raw port scanning and block/allow network traffic routing definitions occur at this physical gateway."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Layer 2: Data Platform (Ingestion)"
    p = tf.add_paragraph()
    p.text = "Tools: Apache NiFi / Atlas, MySQL, FastAPI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: Serves as the primary REST entry-point. Safely validates incoming packets and stores critical metadata into a local relational database for querying."
    p.level = 1
    
    # ----------------------------------------------------
    # Slide 4: L3 & L6 Details
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    body = shapes.placeholders[1]
    
    title.text = "Layer 3 & Layer 6: Semantics & Actions"
    tf = body.text_frame
    
    tf.text = "Layer 3: Knowledge Base (Semantic Memory)"
    p = tf.add_paragraph()
    p.text = "Tools: Qdrant Vector Database"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: Converts standard log text into N-Dimensional arrays, providing a 'Brain' for the AI to dynamically compare if a new threat looks similar to an old one."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Layer 6: Orchestration (Automated Defense)"
    p = tf.add_paragraph()
    p.text = "Tools: n8n Automation Engine"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: Can be triggered programmatically to webhook Discord/Slack channels, block Firewall IPs, or alert Incident response nodes externally."
    p.level = 1
    
    # ----------------------------------------------------
    # Slide 5: L4, L5 & L7 Details
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    body = shapes.placeholders[1]
    
    title.text = "Layer 4/5 & Layer 7: Agents & GUI"
    tf = body.text_frame
    
    tf.text = "Layer 4/5: AI Matrix (Intelligence)"
    p = tf.add_paragraph()
    p.text = "Tools: Mistral / Google Gemini / CrewAI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: Directly parses the analyst's normal English chat inquiries, fetches specific data scopes from MySQL, and synthesizes Threat Analytics into a final PDF."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Layer 7: Agent Web (Unified Visual Plane)"
    p = tf.add_paragraph()
    p.text = "Tools: React, TailwindCSS, Recharts"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Action: The Dashboard surface. Actively polls telemetry, populates dynamic statistics grids, handles User Administration, and offers direct LLM interactions."
    p.level = 1

    # ----------------------------------------------------
    # Slide 6: Summary
    # ----------------------------------------------------
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title = shapes.title
    body = shapes.placeholders[1]
    
    title.text = "Visual Interface Outputs"
    tf = body.text_frame
    tf.text = "When properly connected, the UI intelligently displays:"
    
    p = tf.add_paragraph()
    p.text = "Traffic Real-Time Flow: Interactive AreaChart updating every 2 seconds with aggregated packet drops by severity."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Bot Interactions: The secure chat interface permits on-the-fly 'Give me a report for last 10 minutes' directives returning custom URLs to dynamically crafted Executive PDFs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero-Trust Settings: The administrator gear natively permits Dynamic Theming (Light/Dark mode) alongside Local Administration credentials."
    p.level = 1

    # Save presentation
    filename = 'BrahMos_SOC_Architecture.pptx'
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == '__main__':
    create_brahmos_presentation()
